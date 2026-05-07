"""Generate a PowerPoint deck from the tutorial markdown files.

Each lesson becomes one content slide:
  Title — first H1 from the markdown.
  Bullets — Goal, Recommended dataset, Walkthrough (condensed), Verify.
  Speaker notes — link back to the lesson file in the repo.

Output: slides/xnat-tutorials.pptx
"""

import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

REPO = Path("/Users/james/projects/development/xnat_tutorial_datasets")
OUT = Path("slides/xnat-tutorials.pptx")
HOSTED_DEMO = "https://demo.pro.xnatworks.io"

LEVEL_FOLDERS = ["intro", "intermediate", "advanced", "admin"]
LEVEL_TITLES = {
    "intro": "Intro — UI-only, 5-15 min, no terminal",
    "intermediate": "Intermediate — one tool or concept, 15-30 min",
    "advanced": "Advanced — multi-step pipelines, 30 min+",
    "admin": "Admin — site configuration",
}


def parse_lesson(md_path: Path):
    text = md_path.read_text()
    title = re.search(r"^#\s+(.+)$", text, re.M)
    title = title.group(1).strip() if title else md_path.stem

    meta_line = re.search(r"^\*\*Level:\*\*.+$", text, re.M)
    meta = meta_line.group(0) if meta_line else ""
    dataset = re.search(r"Recommended dataset:\*\*\s*(.+?)\s*$", meta, re.M)
    dataset = dataset.group(1).strip() if dataset else ""

    goal = re.search(r"##\s+Goal\s*\n+([^#]+?)(?=\n##|\Z)", text, re.S)
    goal = goal.group(1).strip().split("\n\n")[0] if goal else ""
    goal = re.sub(r"\s+", " ", goal)

    # Walkthrough: pick first 5 numbered steps
    walk = re.search(
        r"##\s+Walkthrough[^\n]*\n+([^#]+?)(?=\n##|\Z)", text, re.S
    )
    steps = []
    if walk:
        for m in re.finditer(r"^\s*\d+\.\s+(.+?)$", walk.group(1), re.M):
            line = m.group(1).strip()
            line = re.sub(r"\[([^\]]+)\]\([^\)]+\)", r"\1", line)  # de-link
            line = re.sub(r"`([^`]+)`", r"\1", line)
            line = re.sub(r"\*\*([^*]+)\*\*", r"\1", line)
            steps.append(line)
            if len(steps) == 5:
                break

    verify = re.search(r"##\s+Verify\s*\n+([^#]+?)(?=\n##|\Z)", text, re.S)
    verify = verify.group(1).strip().split("\n\n")[0] if verify else ""
    verify = re.sub(r"\s+", " ", verify)
    verify = re.sub(r"`([^`]+)`", r"\1", verify)

    return {
        "title": title,
        "dataset": dataset,
        "goal": goal,
        "steps": steps,
        "verify": verify,
        "rel_path": md_path.relative_to(REPO).as_posix(),
    }


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle
    return slide


def add_section_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[2] if len(prs.slide_layouts) > 2 else prs.slide_layouts[0])
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle
    return slide


def add_content_slide(prs, title, bullets, notes=""):
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    tf.word_wrap = True

    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = b["text"]
        p.level = b.get("level", 0)
        for run in p.runs:
            run.font.size = Pt(b.get("size", 18))

    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. Title slide
    add_title_slide(
        prs,
        "XNAT Tutorials",
        f"Pick-and-mix lesson catalog\n{HOSTED_DEMO}",
    )

    # 2. Overview slide
    add_content_slide(
        prs,
        "Lesson catalog",
        [
            {"text": "9 intro lessons — UI-only, 5-15 min, no terminal", "size": 22},
            {"text": "9 intermediate lessons — one tool or concept, 15-30 min", "size": 22},
            {"text": "5 advanced lessons — multi-step pipelines, 30+ min", "size": 22},
            {"text": "3 admin lessons — site configuration", "size": 22},
            {"text": "Plus glossary, REST cheatsheet, workshop agendas, dataset cards", "size": 18, "level": 1},
        ],
        notes=(
            "Repo: github.com/xnatworks/xnat_tutorial_datasets. "
            "Hosted demo: " + HOSTED_DEMO
        ),
    )

    # 3. How to use
    add_content_slide(
        prs,
        "How to use",
        [
            {"text": "Each intro / intermediate lesson answers four checkpoints", "size": 20},
            {"text": "what changes in XNAT", "level": 1, "size": 18},
            {"text": "where to look", "level": 1, "size": 18},
            {"text": "exact success condition", "level": 1, "size": 18},
            {"text": "where to look first if it fails", "level": 1, "size": 18},
            {"text": "Advanced lessons drop the scaffold for readability", "size": 20},
            {"text": "Datasets are referenced by id and have one-page reference cards", "size": 20},
        ],
    )

    # 4-N. Per-level sections + lesson slides
    for level in LEVEL_FOLDERS:
        folder = REPO / "tutorials" / level
        if not folder.is_dir():
            continue
        files = sorted(folder.glob("*.md"))
        if not files:
            continue

        add_section_slide(prs, level.capitalize(), LEVEL_TITLES[level])

        for md in files:
            data = parse_lesson(md)
            bullets = []
            if data["dataset"]:
                bullets.append({"text": f"Dataset: {data['dataset']}", "size": 18, "level": 1})
            if data["goal"]:
                bullets.append({"text": f"Goal: {data['goal']}", "size": 20})
            for i, step in enumerate(data["steps"], 1):
                bullets.append({"text": f"{i}. {step}", "size": 16, "level": 1})
            if data["verify"]:
                bullets.append({"text": f"Verify: {data['verify']}", "size": 18})

            notes = f"Source: {data['rel_path']}"
            add_content_slide(prs, data["title"], bullets, notes=notes)

    # Closing
    add_content_slide(
        prs,
        "Where to find this",
        [
            {"text": "Repo: github.com/xnatworks/xnat_tutorial_datasets", "size": 22},
            {"text": "Hosted demo: " + HOSTED_DEMO, "size": 22},
            {"text": "Datasets reference: DATASETS.md", "size": 22},
            {"text": "Workshop agendas: tutorials/reference/workshop-agendas.md", "size": 22},
            {"text": "REST cheatsheet: tutorials/reference/rest-cheatsheet.md", "size": 22},
        ],
    )

    # XNAT project resources
    add_content_slide(
        prs,
        "XNAT project resources",
        [
            {"text": "xnat.org — project home", "size": 22},
            {"text": "xnat.org/download/ — download XNAT", "size": 22},
            {"text": "wiki.xnat.org — admin and user wiki", "size": 22},
            {"text": "xnat.readthedocs.io/en/stable/ — Read the Docs", "size": 22},
        ],
    )

    prs.save(str(OUT))
    print(f"Wrote {OUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
