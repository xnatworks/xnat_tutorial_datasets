"""Workflow & Capabilities deck — one slide per pipeline stage.

Mirrors the six sections of the XNAT Workflow & Capabilities table
(INGEST / CURATE / ANNOTATE / COMPUTE / ANALYZE / PUBLISH). Each
slide lists the section's capability bullets and points at the
tutorials that demonstrate them.

Output: slides/xnat-workflow-capabilities.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

REPO = Path(__file__).resolve().parent.parent
OUT = REPO / "slides" / "xnat-workflow-capabilities.pptx"

SECTIONS = [
    {
        "title": "Ingest",
        "tagline": "Capture imaging data into XNAT — scanner, PACS, browser, or REST",
        "bullets": [
            "DICOM C-STORE receiver: scanners send straight to the prearchive",
            "C-MOVE / Q-Retrieve from PACS via XNAT's DICOM SCU",
            "Browser upload, XNAT Desktop Client, REST PUT, and ZIP upload",
            "Site- and project-level DICOM anonymisation scripts (CTP rules)",
            "Pixel-level PHI removal (e.g. dicom-deid-ocr container)",
            "Optional de-facing on structural MR before archive",
            "Modality-agnostic: MR, CT, PET, US, OCT, micro-CT/MR, more",
            "Native handling of DICOM, NIfTI, BIDS, TIFF, vendor formats",
        ],
        "tutorials": [
            "intro/03 — DICOM import and archive",
            "intermediate/02 — De-identification",
        ],
    },
    {
        "title": "Curate",
        "tagline": "Stage, review, and harmonise before data lands in the archive",
        "bullets": [
            "Automated QC as session assessors (MRIQC, custom containers)",
            "Manual review with notes per scan / session, role-based gating",
            "Prearchive staging — reroute, edit, archive, or reject pre-commit",
            "Series-importer rules harmonise scan-type labels at archive time",
            "Project- and cohort-level dashboards over the data table",
            "Custom forms and dynamic data types for clinical / genomic metadata",
            "Format conversion through containers (dcm2niix, dcm2bids, etc.)",
        ],
        "tutorials": [
            "intro/04 — Read DICOM headers",
            "intermediate/01 — DICOM to NIfTI with dcm2niix",
            "intermediate/06 — MRIQC assessor",
            "intermediate/07 — Dynamic data type",
        ],
    },
    {
        "title": "Annotate",
        "tagline": "Mark up images and capture structured reads inside XNAT",
        "bullets": [
            "OHIF Viewer plugin for in-browser DICOM review",
            "ROI / contour / measurement tools, persisted as ROI Collection assessors",
            "AI-assisted segmentation (DEXTR3D, Deepgrow, TotalSegmentator)",
            "Custom-form templates for structured reads, ratings, and QC reviews",
            "Rapid reader module for high-throughput case review",
            "Blinded reader mode with assignment list and audit trail",
            "Radreport.org plugin for structured reporting",
        ],
        "tutorials": [
            "intro/05 — Open a scan in the viewer",
            "intro/09 — Atlases and overlays (Open with Workbench)",
            "intermediate/08 — Custom forms",
            "advanced/05 — TotalSegmentator (AI-assisted RTStruct)",
        ],
    },
    {
        "title": "Compute",
        "tagline": "Run reproducible processing as scan / session / project containers",
        "bullets": [
            "Container Service: register, enable, and launch wrappers per project",
            "HPC submission via Slurm, Torque, PBS, or LSF adapters",
            "Batch processing dashboard across projects and pipelines",
            "Live container logs (stdout / stderr, Tensorboard) in command history",
            "Event Service triggers — auto-run on archive, assessor create, etc.",
        ],
        "tutorials": [
            "intermediate/05 — Container launch basics",
            "admin/01 — Event Service setup",
            "admin/02 — Command enablement",
        ],
    },
    {
        "title": "Analyze",
        "tagline": "Notebook and script access to every XNAT object",
        "bullets": [
            "JupyterHub plugin — notebooks scoped per project, no data egress",
            "NVIDIA Clara / MONAI Label workflows for interactive ML",
            "Group-Level loader for nnU-Net-style train / test / val datasets",
            "xnatpy / xnat-tools / Groovy scripting against the XAPI",
            "Local viewer integrations — 3D Slicer (XNATSlicer), MITK, CapTk",
        ],
        "tutorials": [
            "advanced/01 — Complete BIDS workflow",
            "advanced/02 — nnU-Net dataset builder",
            "advanced/04 — FitLins group analysis",
        ],
    },
    {
        "title": "Publish",
        "tagline": "Share scoped to project, instance, federation, or public registry",
        "bullets": [
            "Cross-project sharing via shared subjects and shared experiments",
            "Cross-instance sharing via XNAT-to-XNAT replication and export",
            "Federated query and dataset assembly across linked XNAT sites",
            "Curated push to TCIA / OpenNeuro via export wrappers",
            "Curated push to Imaging Data Commons / Brainlife.io via containers",
        ],
        "tutorials": [
            "(no dedicated lesson yet — sharing / export to be added)",
        ],
    },
]


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle
    return slide


def add_section_slide(prs, section):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = section["title"]

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = section["tagline"]
    for run in p.runs:
        run.font.size = Pt(20)
        run.font.italic = True

    for bullet in section["bullets"]:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 1
        for run in p.runs:
            run.font.size = Pt(18)

    p = tf.add_paragraph()
    p.text = "Tutorials demonstrating this stage:"
    for run in p.runs:
        run.font.size = Pt(16)
        run.font.bold = True

    for t in section["tutorials"]:
        p = tf.add_paragraph()
        p.text = t
        p.level = 1
        for run in p.runs:
            run.font.size = Pt(15)

    notes = (
        f"{section['title']} stage of the XNAT workflow. "
        f"Capability bullets are taken verbatim from the XNAT Workflow & Capabilities slide. "
        f"Tutorial pointers are slugs under tutorials/ in this repo."
    )
    slide.notes_slide.notes_text_frame.text = notes
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "XNAT Workflow & Capabilities",
        "Six stages: Ingest → Curate → Annotate → Compute → Analyze → Publish",
    )
    for section in SECTIONS:
        add_section_slide(prs, section)

    prs.save(str(OUT))
    print(f"Wrote {OUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
